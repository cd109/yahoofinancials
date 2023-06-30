import openai
import json
import argparse
import time
import random
from pptx import Presentation
from pptx.util import Cm, Pt
import requests, uuid, json


def get_translator_TEXT_MS( S_FROM , S_TO , text ):
  # Add your key and endpoint
  key = "001ca37c284e41e7b0da6e42deea2bce"
  endpoint = "https://api.cognitive.microsofttranslator.com"

  # location, also known as region.
  # required if you're using a multi-service or regional (not global) resource. It can be found in the Azure portal on the Keys and Endpoint page.
  location = "japaneast"

  path = '/translate'
  constructed_url = endpoint + path

  params = {
    'api-version': '3.0',
    'from': S_FROM,
    'to': S_TO
  }

  headers = {
    'Ocp-Apim-Subscription-Key': key,
    # location required if you're using a multi-service or regional (not global) resource.
    'Ocp-Apim-Subscription-Region': location,
    'Content-type': 'application/json',
    'X-ClientTraceId': str(uuid.uuid4())
  }

  # You can pass more than one object in body.
  body = [{
    'text': text
  }]

  request  = requests.post(constructed_url, params=params, headers=headers, json=body)
  response = request.json()
  text_zh  = response[0]['translations'][0]['text'].strip()
  return text_zh #print(text_zh)

def robot_print(text):
    for char in text:
        print(char, end="", flush=True)
        time.sleep(random.randrange(1, 5) / 100.0)
    print("\r")


def chatppt(topic: str, pages: int, api_key: str, language: str):
    language_map = {"cn": "Chinese", "en": "English"}
    language = language_map[language]

    output_format = {
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                # "subtitle": "subtitle for page 1",
                "content": [
                    {
                        "title": "title for bullet 1",
                        "desctription": "detail for bullet 1",
                    },
                    {
                        "title": "title for bullet 2",
                        "desctription": "detail for bullet 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                # "subtitle": "subtitle for page 2",
                "content": [
                    {
                        "title": "title for bullet 1",
                        "desctription": "detail for bullet 1",
                    },
                    {
                        "title": "title for bullet 2",
                        "desctription": "detail for bullet 2",
                    },
                ],
            },
        ],
    }

    messages = [
        {
            "role": "user",
            "content": f"I'm going to prepare a presentation about {topic}, please help to outline detailed about this topic, output with JSON language with follow in format {output_format}, please help to generate {pages} pages, the bullet for each as much as possible, please only return JSON format and use double quotes, please return the content in {language}",
        },
    ]

    robot_print(f"I'm working hard to generate your PPT about {topic}.")
    robot_print("It may takes about a few minutes.")
    robot_print(f"Your PPT will be generated in {language}")

    openai.api_key = api_key
    completion = openai.ChatCompletion.create(model="gpt-3.5-turbo-16k", messages=messages)
    try:
        content = completion.choices[0].message.content
        # just replace ' to " is not a good soluation
        # print(content)
        content = json.loads(content.strip())
        return content
    except Exception as e:
        print("I'm a PPT assistant, your PPT generate failed, please retry later..")
        exit(1)
        # raise Exception("I'm PPT generate assistant, some error happened, please retry")


def generate_ppt(content: str, template=None):
    ppt = Presentation()
    if template:
        ppt = Presentation(template)

    # Creating slide layout
    first_slide_layout = ppt.slide_layouts[0]

    # """ Ref for slide types:
    # 0 ->  title and subtitle
    # 1 ->  title and content
    # 2 ->  section header
    # 3 ->  two content
    # 4 ->  Comparison
    # 5 ->  Title only
    # 6 ->  Blank
    # 7 ->  Content with caption
    # 8 ->  Pic with caption
    # """

    slide = ppt.slides.add_slide(first_slide_layout)
    slide.shapes.title.text =  get_translator_TEXT_MS('en' , 'zh-TW' , content.get("title", "") ) # content.get("title", "")
    slide.placeholders[1].text = "Generate by ChatPPT"

    pages = content.get("pages", [])
    robot_print(f"Your PPT have {len(pages)} pages.")
    for i, page in enumerate(pages):
        page_title = page.get("title", "")
        robot_print(f"page {i+1}: {page_title}")
        bullet_layout = ppt.slide_layouts[1]
        bullet_slide  = ppt.slides.add_slide(bullet_layout)
        bullet_spahe  = bullet_slide.shapes
        bullet_slide.shapes.title.text = get_translator_TEXT_MS('en' , 'zh-TW' , page_title )  # page_title
        
        body_shape = bullet_spahe.placeholders[1]
        for bullet in page.get("content", []):
            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = get_translator_TEXT_MS('en' , 'zh-TW' , bullet.get("title", "") ) # bullet.get("title", "")
            paragraph.font.size = Pt(24)
            paragraph.font.bold = True
            paragraph.level = 1

            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = get_translator_TEXT_MS('en' , 'zh-TW' , bullet.get("description", "") ) # bullet.get("description", "")
            paragraph.font.size = Pt(15)
            paragraph.level = 2

    ppt_name = content.get("title", "")
    ppt_name = f"{ppt_name}.pptx"
    ppt.save(ppt_name)
    robot_print("Generate done, enjoy!")
    robot_print(f"Your PPT: {ppt_name}")


def main():
    robot_print("Hi, I am your PPT assistant.")
    robot_print("I am powered by ChatGPT")
    # robot_print("If you have any issue, please contact hui_mi@dell.com")
    ppt_content = chatppt("王安石變法實際案例", 5 , "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", "en")
    generate_ppt(ppt_content)


if __name__ == "__main__":
    main()
